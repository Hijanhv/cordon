"""Cardon Protocol pitch deck — clean, minimal, no overlapping elements."""

import io
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ───────────────────────────────────────────────────────────────────
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLACK   = RGBColor(0x0D, 0x0D, 0x0D)
GRAY    = RGBColor(0x6B, 0x72, 0x80)
LGRAY   = RGBColor(0xF3, 0xF4, 0xF6)
BORDER  = RGBColor(0xE5, 0xE7, 0xEB)
GREEN   = RGBColor(0x16, 0xA3, 0x4A)
BLUE    = RGBColor(0x22, 0x56, 0xD8)

W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Primitives ────────────────────────────────────────────────────────────────

def bg(slide):
    s = slide.shapes.add_shape(1, 0, 0, W, H)
    s.fill.solid(); s.fill.fore_color.rgb = WHITE
    s.line.fill.background()

def rect(slide, l, t, w, h, fill=None, border=None, bw=Pt(0.5)):
    s = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        s.fill.solid(); s.fill.fore_color.rgb = fill
    else:
        s.fill.background()
    if border:
        s.line.color.rgb = border; s.line.width = bw
    else:
        s.line.fill.background()
    return s

def txt(slide, text, l, t, w, h,
        size=14, bold=False, color=BLACK,
        align=PP_ALIGN.LEFT, wrap=True):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold
    r.font.color.rgb = color; r.font.name = "Helvetica Neue"
    return tb

def header(slide, slide_label):
    """Top bar: brand left, slide label right."""
    rect(slide, 0, 0, W, Inches(0.55), fill=BLACK)
    txt(slide, "Cardon Protocol",
        Inches(0.4), Inches(0.08), Inches(5), Inches(0.4),
        size=15, bold=True, color=WHITE)
    txt(slide, slide_label,
        Inches(0), Inches(0.08), W - Inches(0.4), Inches(0.4),
        size=11, color=RGBColor(0xA0,0xA0,0xA0), align=PP_ALIGN.RIGHT)

def footer(slide, page):
    """Bottom bar: URL left, page number right."""
    rect(slide, 0, H - Inches(0.35), W, Inches(0.35), fill=LGRAY)
    txt(slide, "cardon-protocol.vercel.app",
        Inches(0.4), H - Inches(0.32), Inches(6), Inches(0.3),
        size=9, color=GRAY)
    txt(slide, str(page),
        Inches(0), H - Inches(0.32), W - Inches(0.4), Inches(0.3),
        size=9, color=GRAY, align=PP_ALIGN.RIGHT)

def divider(slide, t):
    rect(slide, Inches(0.5), t, W - Inches(1.0), Inches(0.015), fill=BORDER)

# Content area: x=0.5, y starts after header (0.55) + padding (0.2) = 0.75
# Content bottom: above footer 7.5 - 0.35 - 0.15 = 7.0  → height = 6.25"

# ── Terminal screenshot PNG ────────────────────────────────────────────────────

def make_terminal_png() -> io.BytesIO:
    scale = 2
    W_px, H_px = 1060 * scale, 460 * scale
    sc = scale

    BG    = (22,  22,  35)
    BAR   = (40,  40,  56)
    PASS  = (100, 200, 120)
    DIM   = (100, 105, 130)
    WHT   = (220, 225, 245)
    BLU   = (110, 170, 250)

    img = Image.new("RGB", (W_px, H_px), BG)
    d   = ImageDraw.Draw(img)

    d.rectangle([0, 0, W_px, 38 * sc], fill=BAR)
    for cx, col in [(22*sc,(220,90,100)),(42*sc,(240,190,90)),(62*sc,(100,200,120))]:
        d.ellipse([cx-6*sc, 11*sc, cx+6*sc, 23*sc], fill=col)

    try:
        mono = ImageFont.truetype("/System/Library/Fonts/Menlo.ttc", 14*sc)
        mono_sm = ImageFont.truetype("/System/Library/Fonts/Menlo.ttc", 12*sc)
        sans = ImageFont.truetype("/Library/Fonts/Arial.ttf", 12*sc)
    except Exception:
        mono = mono_sm = sans = ImageFont.load_default()

    title = "zsh — cordon — anchor test --provider.cluster devnet"
    bb = d.textbbox((0,0), title, font=sans)
    d.text(((W_px - (bb[2]-bb[0]))//2, 12*sc), title, font=sans, fill=DIM)

    d.text((20*sc, 52*sc),
           "$ anchor test --provider.cluster devnet --skip-local-validator",
           font=mono_sm, fill=DIM)
    d.text((20*sc, 76*sc), "  cardon", font=mono, fill=WHT)

    tests = [
        ("registers an agent and seeds an initial policy",           "(4218ms)"),
        ("auto-approves a small tx to an allowed program",           "(3594ms)"),
        ("rejects an auto submission over the HITL threshold",       " (939ms)"),
        ("rejects an auto submission to a non-allowlisted program",  " (848ms)"),
        ("queues a high-value tx for human review and approves it",  "(7564ms)"),
        ("rejects a pending tx",                                     "(3499ms)"),
        ("blocks submissions when the agent is disabled",            "(4142ms)"),
        ("rejects registration when hitl_threshold > max_per_tx",   " (768ms)"),
        ("rejects a submission that would breach daily volume cap",  "(3740ms)"),
    ]

    y = 102 * sc
    row = 26 * sc
    for name, timing in tests:
        d.text((20*sc, y),  "    ✔", font=mono, fill=PASS)
        d.text((74*sc, y),  name,    font=mono, fill=WHT)
        bb2 = d.textbbox((0,0), timing, font=mono_sm)
        d.text((W_px - (bb2[2]-bb2[0]) - 20*sc, y), timing, font=mono_sm, fill=DIM)
        y += row

    y += 6*sc
    d.line([(20*sc, y), (W_px-20*sc, y)], fill=BAR, width=sc)
    y += 14*sc
    d.text((20*sc, y), "  9 passing", font=mono, fill=PASS)
    d.text((168*sc, y), "(33s)", font=mono_sm, fill=DIM)
    y += 28*sc
    d.line([(20*sc, y), (W_px-20*sc, y)], fill=BAR, width=sc)
    y += 12*sc
    d.text((20*sc, y), "Program ID (devnet):", font=mono_sm, fill=DIM)
    d.text((230*sc, y), "5Zy73PLH2hpNvsCvyeqpn2tKjHBXNihjX7UbeAZ41tQc",
           font=mono_sm, fill=BLU)

    out = img.resize((W_px//scale, H_px//scale), Image.LANCZOS)
    buf = io.BytesIO(); out.save(buf, "PNG"); buf.seek(0)
    return buf


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
bg(s1)

# Full-width black top band — brand name lives here
rect(s1, 0, 0, W, Inches(1.1), fill=BLACK)
txt(s1, "Cardon Protocol",
    Inches(0.55), Inches(0.22), Inches(8), Inches(0.7),
    size=32, bold=True, color=WHITE)
txt(s1, "Turbin3 Builders Cohort — Capstone",
    Inches(0), Inches(0.28), W - Inches(0.5), Inches(0.55),
    size=12, color=GRAY, align=PP_ALIGN.RIGHT)

# Tagline block — vertically centred in remaining space
txt(s1,
    "An on-chain transaction firewall\nfor Solana AI agents.",
    Inches(0.55), Inches(1.55), Inches(10), Inches(2.2),
    size=38, bold=True, color=BLACK)

divider(s1, Inches(3.95))

txt(s1,
    "Cardon sits between an autonomous agent and the network. "
    "Every transaction is simulated, checked against a policy stored in an "
    "Anchor program, and either auto-approved, queued for human review, or rejected. "
    "Because the policy lives on-chain, no operator can quietly raise a spend cap.",
    Inches(0.55), Inches(4.1), Inches(12.2), Inches(1.6),
    size=14, color=GRAY)

divider(s1, Inches(5.9))

txt(s1, "Live demo",         Inches(0.55), Inches(6.05), Inches(2), Inches(0.3), size=10, bold=True, color=GRAY)
txt(s1, "cardon-protocol.vercel.app",
    Inches(0.55), Inches(6.38), Inches(5), Inches(0.35), size=13, bold=False, color=BLUE)

txt(s1, "Program ID (devnet)",  Inches(5.5), Inches(6.05), Inches(7), Inches(0.3), size=10, bold=True, color=GRAY)
txt(s1, "5Zy73PLH2hpNvsCvyeqpn2tKjHBXNihjX7UbeAZ41tQc",
    Inches(5.5), Inches(6.38), Inches(7.4), Inches(0.35), size=11, color=BLACK)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ═══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
bg(s2); header(s2, "The Problem"); footer(s2, 2)

txt(s2, "The Problem",
    Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=BLACK)

divider(s2, Inches(1.45))

problems = [
    ("No firewall between intent and the network",
     "Autonomous agents sign and broadcast transactions on their own. A single bad prompt, "
     "jailbreak, or model regression can drain an entire hot wallet — instantly and irreversibly."),
    ("Centralized policy is silently mutable",
     "Existing products run policy as centralized middleware. Anyone with code access can "
     "quietly raise a spend cap or remove a kill switch with no audit trail visible to token holders."),
    ("No programmable human-in-the-loop",
     "There is no standard on-chain primitive to pause high-value transactions and route "
     "them to a human approver before they hit the network."),
]

y = Inches(1.6)
for title, body in problems:
    rect(s2, Inches(0.5), y, Inches(0.06), Inches(0.3), fill=BLACK)
    txt(s2, title,
        Inches(0.75), y - Inches(0.04), Inches(12.0), Inches(0.4),
        size=15, bold=True, color=BLACK)
    txt(s2, body,
        Inches(0.75), y + Inches(0.36), Inches(12.0), Inches(0.7),
        size=12, color=GRAY)
    y += Inches(1.55)

# callout box at bottom
rect(s2, Inches(0.5), Inches(6.18), Inches(12.33), Inches(0.6),
     fill=LGRAY, border=BORDER)
txt(s2,
    "\"One bad prompt can drain a hot wallet — and there is currently no firewall between an agent's intent and the network.\"",
    Inches(0.7), Inches(6.25), Inches(11.9), Inches(0.48),
    size=12, color=GRAY, align=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — The Solution + Benefits
# ═══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
bg(s3); header(s3, "The Solution"); footer(s3, 3)

txt(s3, "The Solution",
    Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=BLACK)

divider(s3, Inches(1.45))

# 3-step flow cards
flow = [
    ("1  Simulate",
     "Every transaction is run through RPC simulation before anything is signed."),
    ("2  Check Policy",
     "Checked on-chain: per-tx cap, program allowlist, daily volume ceiling, kill switch."),
    ("3  Decide",
     "Under threshold → auto-approve.  Over threshold → queue for human.  Cap hit → reject."),
]

bx = Inches(0.5)
bw = Inches(3.9)
gap = Inches(0.27)
for title, body in flow:
    rect(s3, bx, Inches(1.6), bw, Inches(1.9), fill=LGRAY, border=BORDER)
    txt(s3, title,
        bx + Inches(0.2), Inches(1.72), bw - Inches(0.35), Inches(0.45),
        size=14, bold=True, color=BLACK)
    txt(s3, body,
        bx + Inches(0.2), Inches(2.22), bw - Inches(0.35), Inches(1.1),
        size=12, color=GRAY)
    bx += bw + gap

divider(s3, Inches(3.68))

# Benefits — 2 × 2 grid
txt(s3, "Key Benefits",
    Inches(0.5), Inches(3.82), Inches(6), Inches(0.4),
    size=14, bold=True, color=BLACK)

benefits = [
    ("Immutable authority",
     "Policy authority can be a Squads multisig or DAO — rotating rules is an on-chain transaction, not a config push."),
    ("Full audit trail",
     "Every approval, rejection, and policy change is emitted as a log event, indexed off-chain via Helius."),
    ("Human-in-the-loop built in",
     "High-value transactions are automatically queued for human review before broadcasting."),
    ("Zero-trust kill switch",
     "set_enabled can freeze an agent instantly — no operator coordination or code deploy required."),
]

cols  = [Inches(0.5), Inches(6.95)]
row_y = [Inches(4.3), Inches(5.5)]
for i, (title, body) in enumerate(benefits):
    cx = cols[i % 2]
    cy = row_y[i // 2]
    rect(s3, cx, cy + Inches(0.08), Inches(0.07), Inches(0.07), fill=GREEN)
    txt(s3, title,
        cx + Inches(0.22), cy, Inches(5.8), Inches(0.38),
        size=13, bold=True, color=BLACK)
    txt(s3, body,
        cx + Inches(0.22), cy + Inches(0.38), Inches(5.8), Inches(0.72),
        size=11, color=GRAY)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Tests Passing
# ═══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
bg(s4); header(s4, "Verified on Devnet"); footer(s4, 4)

txt(s4, "All 9 Tests Passing on Devnet",
    Inches(0.5), Inches(0.72), Inches(12), Inches(0.65),
    size=28, bold=True, color=BLACK)

divider(s4, Inches(1.45))

# terminal screenshot — full width under the title
term = make_terminal_png()
s4.shapes.add_picture(term, Inches(0.5), Inches(1.6), Inches(8.9), Inches(3.88))

# right info panel — positioned to match the terminal height exactly
PX = Inches(9.65)
PW = Inches(3.18)
py = Inches(1.6)
rh = Inches(0.82)
gap = Inches(0.12)

for label, value, vc, fc in [
    ("NETWORK",     "Solana Devnet",             GREEN, LGRAY),
    ("STATUS",      "9 / 9 passing  ✓",          GREEN, LGRAY),
    ("PROGRAM ID",  "5Zy73PLH2hpN\nvsCvyeqpn2tKjHBXNihjX7UbeAZ41tQc",
                                                  BLACK, LGRAY),
    ("LIVE SITE",   "cardon-protocol\n.vercel.app", BLUE, LGRAY),
]:
    rect(s4, PX, py, PW, rh, fill=fc, border=BORDER)
    txt(s4, label,
        PX + Inches(0.15), py + Inches(0.08), PW - Inches(0.2), Inches(0.22),
        size=8, bold=True, color=GRAY)
    txt(s4, value,
        PX + Inches(0.15), py + Inches(0.3), PW - Inches(0.2), Inches(0.48),
        size=11, bold=True, color=vc)
    py += rh + gap

# ── Save ──────────────────────────────────────────────────────────────────────
out = "/Users/janhavi/Desktop/cordon/cardon-pitch-deck.pptx"
prs.save(out)
print(f"Saved → {out}")
